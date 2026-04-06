import os
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

SUPPORTED = (".pdf", ".txt", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".csv", ".md")

CATEGORIES = ["health", "car", "banking", "home"]

def load_all_categories(data_dir: str):
    """
    Load documents from each subfolder separately.
    Returns dict: { "health": [docs], "car": [docs], ... }
    """
    category_docs  = {}
    category_names = {}

    for category in CATEGORIES:
        folder = os.path.join(data_dir, category)
        if not os.path.exists(folder):
            continue
        docs, names = load_folder(folder, category)
        if docs:
            category_docs[category]  = docs
            category_names[category] = names
            print(f"✓ {category}: {len(docs)} chunks from {len(names)} files")

    return category_docs, category_names


def load_folder(folder: str, category: str):
    documents, doc_names = [], []
    for filename in sorted(os.listdir(folder)):
        ext = os.path.splitext(filename)[1].lower()
        if ext not in SUPPORTED:
            continue
        filepath = os.path.join(folder, filename)
        doc_names.append(filename)
        try:
            text = _extract(filepath, ext)
            if text and text.strip():
                chunks = _split(text, filename, category)
                documents.extend(chunks)
        except Exception as e:
            print(f"✗ {filename}: {e}")
    return documents, doc_names


def _extract(filepath, ext):
    if ext == ".pdf":
        from pypdf import PdfReader
        return "\n".join(p.extract_text() or "" for p in PdfReader(filepath).pages)
    elif ext in (".docx", ".doc"):
        import docx
        d = docx.Document(filepath)
        parts = [p.text for p in d.paragraphs if p.text.strip()]
        for t in d.tables:
            for row in t.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n".join(parts)
    elif ext in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        parts = []
        for s in wb.sheetnames:
            parts.append(f"Sheet: {s}")
            for row in wb[s].iter_rows(values_only=True):
                r = " | ".join(str(c) for c in row if c is not None)
                if r.strip(): parts.append(r)
        return "\n".join(parts)
    elif ext in (".pptx", ".ppt"):
        from pptx import Presentation
        prs = Presentation(filepath)
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)
    elif ext == ".csv":
        import csv
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return "\n".join(" | ".join(r) for r in csv.reader(f))
    elif ext in (".txt", ".md"):
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    return ""


def _split(text, source, category):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=400,
        chunk_overlap=80,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    return [
        Document(page_content=c, metadata={"source": source, "category": category})
        for c in splitter.split_text(text) if c.strip()
    ]
